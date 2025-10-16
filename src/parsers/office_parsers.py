"""
Office 文档解析器模块

提供 Word、Excel、PowerPoint 文档的解析功能
"""

from typing import Dict, Any, List
import logging
import os
from docx import Document
from docx.table import Table
from openpyxl import load_workbook
from pptx import Presentation
from .base import BaseParser, ParseResult

logger = logging.getLogger(__name__)


class DocxParser(BaseParser):
    """
    Word 文档解析器 (.docx)

    支持提取:
    - 正文段落
    - 表格内容
    - 页眉和页脚
    - 文档元数据(作者、标题、创建时间等)
    """

    def __init__(self):
        """初始化 Word 文档解析器"""
        super().__init__(supported_extensions=['.docx'])

    def _parse_impl(self, file_path: str) -> ParseResult:
        """
        解析 Word 文档

        Args:
            file_path: Word 文档路径

        Returns:
            解析结果
        """
        try:
            doc = Document(file_path)

            # 提取正文段落
            content_parts = []
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    content_parts.append(text)

            # 提取表格内容
            for table in doc.tables:
                table_text = self._extract_table_text(table)
                if table_text:
                    content_parts.append(table_text)

            # 提取页眉和页脚
            headers_footers = self._extract_headers_footers(doc)
            if headers_footers:
                content_parts.append(headers_footers)

            # 合并内容
            content = '\n'.join(content_parts)

            # 提取元数据
            metadata = self._extract_metadata(doc, file_path)

            return ParseResult(
                success=True,
                content=content,
                metadata=metadata
            )

        except Exception as e:
            logger.error(f"解析 Word 文档失败 {file_path}: {e}")
            return ParseResult(
                success=False,
                content='',
                error=f'解析 Word 文档失败: {str(e)}'
            )

    def _extract_table_text(self, table: Table) -> str:
        """
        提取表格文本

        Args:
            table: Word 表格对象

        Returns:
            格式化的表格文本
        """
        table_lines = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                text = cell.text.strip()
                cells.append(text if text else '')
            if any(cells):  # 只添加非空行
                table_lines.append(' | '.join(cells))

        return '\n'.join(table_lines)

    def _extract_headers_footers(self, doc: Document) -> str:
        """
        提取页眉和页脚

        Args:
            doc: Word 文档对象

        Returns:
            页眉页脚文本
        """
        parts = []

        for section in doc.sections:
            # 提取页眉
            if section.header:
                header_text = '\n'.join(
                    p.text.strip() for p in section.header.paragraphs
                    if p.text.strip()
                )
                if header_text:
                    parts.append(f"页眉: {header_text}")

            # 提取页脚
            if section.footer:
                footer_text = '\n'.join(
                    p.text.strip() for p in section.footer.paragraphs
                    if p.text.strip()
                )
                if footer_text:
                    parts.append(f"页脚: {footer_text}")

        return '\n'.join(parts)

    def _extract_metadata(self, doc: Document, file_path: str) -> Dict[str, Any]:
        """
        提取文档元数据

        Args:
            doc: Word 文档对象
            file_path: 文件路径

        Returns:
            元数据字典
        """
        core_props = doc.core_properties

        metadata = {
            'file_type': 'docx',
            'author': core_props.author or '',
            'title': core_props.title or '',
            'subject': core_props.subject or '',
            'keywords': core_props.keywords or '',
            'comments': core_props.comments or '',
            'created': core_props.created.isoformat() if core_props.created else '',
            'modified': core_props.modified.isoformat() if core_props.modified else '',
            'last_modified_by': core_props.last_modified_by or '',
            'paragraphs': len(doc.paragraphs),
            'tables': len(doc.tables),
            'sections': len(doc.sections),
            'file_size': os.path.getsize(file_path)
        }

        return metadata


class XlsxParser(BaseParser):
    """
    Excel 文档解析器 (.xlsx, .xls)

    支持提取:
    - 所有工作表内容
    - 单元格文本和数值
    - 单元格公式(转为文本)
    - 合并单元格处理
    - 工作簿属性
    """

    def __init__(self):
        """初始化 Excel 文档解析器"""
        super().__init__(supported_extensions=['.xlsx', '.xls'])

    def _parse_impl(self, file_path: str) -> ParseResult:
        """
        解析 Excel 文档

        Args:
            file_path: Excel 文档路径

        Returns:
            解析结果
        """
        try:
            # data_only=True 将公式转换为值
            workbook = load_workbook(file_path, data_only=True, read_only=True)

            content_parts = []
            sheet_info = []
            total_cells = 0

            # 遍历所有工作表
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_info.append(sheet_name)

                # 添加工作表标题
                content_parts.append(f"\n【工作表: {sheet_name}】")

                # 提取工作表内容
                sheet_content, cell_count = self._extract_sheet_content(sheet)
                if sheet_content:
                    content_parts.append(sheet_content)
                total_cells += cell_count

            # 合并内容
            content = '\n'.join(content_parts)

            # 提取元数据
            metadata = self._extract_metadata(workbook, file_path, sheet_info, total_cells)

            workbook.close()

            return ParseResult(
                success=True,
                content=content,
                metadata=metadata
            )

        except Exception as e:
            logger.error(f"解析 Excel 文档失败 {file_path}: {e}")
            return ParseResult(
                success=False,
                content='',
                error=f'解析 Excel 文档失败: {str(e)}'
            )

    def _extract_sheet_content(self, sheet) -> tuple[str, int]:
        """
        提取工作表内容

        Args:
            sheet: Excel 工作表对象

        Returns:
            (格式化的内容, 单元格数量)
        """
        rows = []
        cell_count = 0

        # 获取有数据的区域
        if sheet.max_row == 0 or sheet.max_column == 0:
            return '', 0

        for row in sheet.iter_rows(values_only=True):
            # 转换单元格值为字符串
            cells = []
            for cell in row:
                if cell is not None:
                    cell_count += 1
                    # 处理不同类型的值
                    if isinstance(cell, (int, float)):
                        cells.append(str(cell))
                    else:
                        cells.append(str(cell).strip())
                else:
                    cells.append('')

            # 只添加非空行
            row_text = ' | '.join(cells)
            if row_text.strip(' |'):
                rows.append(row_text)

        return '\n'.join(rows), cell_count

    def _extract_metadata(
        self,
        workbook,
        file_path: str,
        sheet_names: List[str],
        total_cells: int
    ) -> Dict[str, Any]:
        """
        提取工作簿元数据

        Args:
            workbook: Excel 工作簿对象
            file_path: 文件路径
            sheet_names: 工作表名称列表
            total_cells: 总单元格数

        Returns:
            元数据字典
        """
        props = workbook.properties

        metadata = {
            'file_type': 'xlsx',
            'sheets': len(sheet_names),
            'sheet_names': sheet_names,
            'total_cells': total_cells,
            'title': props.title or '',
            'author': props.creator or '',
            'subject': props.subject or '',
            'keywords': props.keywords or '',
            'comments': props.description or '',
            'created': props.created.isoformat() if props.created else '',
            'modified': props.modified.isoformat() if props.modified else '',
            'last_modified_by': props.lastModifiedBy or '',
            'file_size': os.path.getsize(file_path)
        }

        return metadata


class PptxParser(BaseParser):
    """
    PowerPoint 文档解析器 (.pptx, .ppt)

    支持提取:
    - 所有幻灯片标题
    - 幻灯片正文内容
    - 演讲者备注
    - 幻灯片顺序
    - 演示文稿属性
    """

    def __init__(self):
        """初始化 PowerPoint 文档解析器"""
        super().__init__(supported_extensions=['.pptx', '.ppt'])

    def _parse_impl(self, file_path: str) -> ParseResult:
        """
        解析 PowerPoint 文档

        Args:
            file_path: PowerPoint 文档路径

        Returns:
            解析结果
        """
        try:
            prs = Presentation(file_path)

            content_parts = []
            slide_count = len(prs.slides)

            # 遍历所有幻灯片
            for i, slide in enumerate(prs.slides, 1):
                slide_parts = [f"\n【幻灯片 {i}/{slide_count}】"]

                # 提取标题和正文
                text_content = self._extract_slide_text(slide)
                if text_content:
                    slide_parts.append(text_content)

                # 提取备注
                notes = self._extract_slide_notes(slide)
                if notes:
                    slide_parts.append(f"备注: {notes}")

                content_parts.extend(slide_parts)

            # 合并内容
            content = '\n'.join(content_parts)

            # 提取元数据
            metadata = self._extract_metadata(prs, file_path, slide_count)

            return ParseResult(
                success=True,
                content=content,
                metadata=metadata
            )

        except Exception as e:
            logger.error(f"解析 PowerPoint 文档失败 {file_path}: {e}")
            return ParseResult(
                success=False,
                content='',
                error=f'解析 PowerPoint 文档失败: {str(e)}'
            )

    def _extract_slide_text(self, slide) -> str:
        """
        提取幻灯片文本

        Args:
            slide: PowerPoint 幻灯片对象

        Returns:
            幻灯片文本
        """
        text_parts = []

        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue

            text = shape.text.strip()
            if text:
                # 判断是否为标题
                if shape.is_placeholder:
                    placeholder = shape.placeholder_format
                    if placeholder.type == 1:  # 标题占位符
                        text_parts.append(f"标题: {text}")
                        continue

                text_parts.append(text)

        return '\n'.join(text_parts)

    def _extract_slide_notes(self, slide) -> str:
        """
        提取幻灯片备注

        Args:
            slide: PowerPoint 幻灯片对象

        Returns:
            备注文本
        """
        if not slide.has_notes_slide:
            return ''

        try:
            notes_slide = slide.notes_slide
            if hasattr(notes_slide, 'notes_text_frame'):
                text = notes_slide.notes_text_frame.text.strip()
                return text
        except Exception as e:
            logger.warning(f"提取备注失败: {e}")

        return ''

    def _extract_metadata(
        self,
        prs: Presentation,
        file_path: str,
        slide_count: int
    ) -> Dict[str, Any]:
        """
        提取演示文稿元数据

        Args:
            prs: PowerPoint 演示文稿对象
            file_path: 文件路径
            slide_count: 幻灯片数量

        Returns:
            元数据字典
        """
        core_props = prs.core_properties

        metadata = {
            'file_type': 'pptx',
            'slides': slide_count,
            'title': core_props.title or '',
            'author': core_props.author or '',
            'subject': core_props.subject or '',
            'keywords': core_props.keywords or '',
            'comments': core_props.comments or '',
            'created': core_props.created.isoformat() if core_props.created else '',
            'modified': core_props.modified.isoformat() if core_props.modified else '',
            'last_modified_by': core_props.last_modified_by or '',
            'file_size': os.path.getsize(file_path)
        }

        # 添加幻灯片尺寸信息
        try:
            metadata['slide_width'] = prs.slide_width
            metadata['slide_height'] = prs.slide_height
        except:
            pass

        return metadata
