"""
创建测试样本文件

生成包含中文内容的 Office 文档用于测试
"""

import os
from pathlib import Path
from docx import Document
from docx.shared import Pt, Inches
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt


def create_sample_docx():
    """创建测试 Word 文档"""
    doc = Document()

    # 设置文档属性
    doc.core_properties.author = '测试作者'
    doc.core_properties.title = 'Word 测试文档'
    doc.core_properties.subject = '文档解析测试'
    doc.core_properties.keywords = 'Python, 测试, Word'
    doc.core_properties.comments = '这是一个用于测试的 Word 文档'

    # 添加标题
    doc.add_heading('测试 Word 文档', 0)

    # 添加段落
    doc.add_paragraph('这是第一段内容,包含中文和English混合文本。')
    doc.add_paragraph('第二段测试特殊字符:!@#$%^&*()_+-=[]{}|;:\'",.<>?/')

    # 添加子标题
    doc.add_heading('表格测试', level=1)

    # 添加表格
    table = doc.add_table(rows=4, cols=3)
    table.style = 'Light Grid Accent 1'

    # 表格标题
    headers = ['姓名', '年龄', '职位']
    for i, header in enumerate(headers):
        table.rows[0].cells[i].text = header

    # 表格数据
    data = [
        ['张三', '28', '工程师'],
        ['李四', '32', '设计师'],
        ['王五', '25', '产品经理']
    ]

    for i, row_data in enumerate(data, start=1):
        for j, cell_value in enumerate(row_data):
            table.rows[i].cells[j].text = cell_value

    # 添加列表
    doc.add_heading('列表测试', level=1)
    doc.add_paragraph('功能列表:', style='List Bullet')
    doc.add_paragraph('文本提取', style='List Bullet')
    doc.add_paragraph('表格处理', style='List Bullet')
    doc.add_paragraph('元数据读取', style='List Bullet')

    # 设置页眉
    section = doc.sections[0]
    header = section.header
    header_para = header.paragraphs[0]
    header_para.text = "测试文档页眉"

    # 设置页脚
    footer = section.footer
    footer_para = footer.paragraphs[0]
    footer_para.text = "第 1 页"

    return doc


def create_sample_xlsx():
    """创建测试 Excel 文档"""
    wb = Workbook()

    # 设置工作簿属性
    wb.properties.creator = '测试作者'
    wb.properties.title = 'Excel 测试文档'
    wb.properties.subject = 'Excel 解析测试'
    wb.properties.keywords = 'Python, 测试, Excel'
    wb.properties.description = '这是一个用于测试的 Excel 文档'

    # 第一个工作表 - 员工信息
    ws1 = wb.active
    ws1.title = "员工信息"

    # 添加表头
    headers = ['工号', '姓名', '部门', '职位', '工资']
    ws1.append(headers)

    # 添加数据
    data = [
        ['E001', '张三', '技术部', 'Python工程师', 15000],
        ['E002', '李四', '设计部', 'UI设计师', 12000],
        ['E003', '王五', '产品部', '产品经理', 18000],
        ['E004', '赵六', '技术部', 'Java工程师', 16000],
        ['E005', '钱七', '市场部', '市场专员', 10000],
    ]

    for row in data:
        ws1.append(row)

    # 添加合计行
    ws1.append(['', '', '', '平均工资:', '=AVERAGE(E2:E6)'])

    # 第二个工作表 - 销售数据
    ws2 = wb.create_sheet(title="销售数据")

    # 添加表头
    sales_headers = ['日期', '产品', '数量', '单价', '金额']
    ws2.append(sales_headers)

    # 添加数据
    sales_data = [
        ['2024-01-01', '产品A', 100, 50, '=C2*D2'],
        ['2024-01-02', '产品B', 150, 80, '=C3*D3'],
        ['2024-01-03', '产品C', 200, 60, '=C4*D4'],
        ['2024-01-04', '产品A', 120, 50, '=C5*D5'],
        ['2024-01-05', '产品B', 180, 80, '=C6*D6'],
    ]

    for row in sales_data:
        ws2.append(row)

    # 添加合计
    ws2.append(['', '合计', '=SUM(C2:C6)', '', '=SUM(E2:E6)'])

    # 第三个工作表 - 测试特殊字符
    ws3 = wb.create_sheet(title="特殊字符测试")
    ws3.append(['特殊字符', '描述'])
    ws3.append(['!@#$%', '符号测试'])
    ws3.append(['中文,English', '混合语言'])
    ws3.append(['12345', '数字'])
    ws3.append(['空白单元格测试', None])

    return wb


def create_sample_pptx():
    """创建测试 PowerPoint 文档"""
    prs = Presentation()

    # 设置演示文稿属性
    prs.core_properties.author = '测试作者'
    prs.core_properties.title = 'PowerPoint 测试文档'
    prs.core_properties.subject = 'PPT 解析测试'
    prs.core_properties.keywords = 'Python, 测试, PowerPoint'
    prs.core_properties.comments = '这是一个用于测试的 PowerPoint 文档'

    # 幻灯片 1 - 标题页
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    title1.text = "Windows 搜索工具"
    subtitle1.text = "文档解析引擎测试\n测试作者"

    # 添加备注
    notes_slide1 = slide1.notes_slide
    notes_slide1.notes_text_frame.text = "这是第一张幻灯片的备注,用于测试备注提取功能。"

    # 幻灯片 2 - 内容页
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title2 = slide2.shapes.title
    title2.text = "功能特性"

    # 添加文本框
    body2 = slide2.placeholders[1]
    tf2 = body2.text_frame
    tf2.text = "主要功能:"

    # 添加列表项
    features = [
        "支持多种文档格式",
        "智能内容提取",
        "中文分词与索引",
        "快速全文搜索"
    ]

    for feature in features:
        p = tf2.add_paragraph()
        p.text = feature
        p.level = 1

    # 添加备注
    notes_slide2 = slide2.notes_slide
    notes_slide2.notes_text_frame.text = "功能特性页面,展示系统主要能力。"

    # 幻灯片 3 - 技术架构
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = "技术架构"

    body3 = slide3.placeholders[1]
    tf3 = body3.text_frame
    tf3.text = "核心组件:"

    components = [
        "文档解析引擎 (python-docx, openpyxl, python-pptx)",
        "分词索引模块 (jieba)",
        "数据库存储 (SQLite)",
        "图形界面 (PyQt6)"
    ]

    for component in components:
        p = tf3.add_paragraph()
        p.text = component
        p.level = 1

    # 添加备注
    notes_slide3 = slide3.notes_slide
    notes_slide3.notes_text_frame.text = "技术架构采用模块化设计,便于扩展和维护。"

    # 幻灯片 4 - 特殊字符测试
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    title4.text = "特殊字符测试"

    body4 = slide4.placeholders[1]
    tf4 = body4.text_frame
    tf4.text = "特殊字符: !@#$%^&*()_+-=[]{}|"

    p = tf4.add_paragraph()
    p.text = "中英混合: Hello 你好 World 世界"
    p.level = 0

    p = tf4.add_paragraph()
    p.text = "数字测试: 12345 67890"
    p.level = 0

    # 幻灯片 5 - 总结
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    title5.text = "总结"

    body5 = slide5.placeholders[1]
    tf5 = body5.text_frame
    tf5.text = "测试要点:"

    summary = [
        "文本提取测试通过",
        "表格处理测试通过",
        "元数据读取测试通过",
        "中文支持测试通过"
    ]

    for item in summary:
        p = tf5.add_paragraph()
        p.text = item
        p.level = 1

    # 添加备注
    notes_slide5 = slide5.notes_slide
    notes_slide5.notes_text_frame.text = "总结页面,确认所有测试项目均已通过。"

    return prs


def main():
    """主函数:生成所有测试样本文件"""
    # 确定输出目录
    script_dir = Path(__file__).parent
    output_dir = script_dir / 'sample_files'
    output_dir.mkdir(exist_ok=True)

    print("开始创建测试样本文件...")

    # 创建 Word 文档
    print("创建 Word 文档...")
    docx_file = output_dir / 'sample.docx'
    doc = create_sample_docx()
    doc.save(docx_file)
    print(f"[OK] 已创建: {docx_file}")

    # 创建 Excel 文档
    print("创建 Excel 文档...")
    xlsx_file = output_dir / 'sample.xlsx'
    wb = create_sample_xlsx()
    wb.save(xlsx_file)
    print(f"[OK] 已创建: {xlsx_file}")

    # 创建 PowerPoint 文档
    print("创建 PowerPoint 文档...")
    pptx_file = output_dir / 'sample.pptx'
    prs = create_sample_pptx()
    prs.save(pptx_file)
    print(f"[OK] 已创建: {pptx_file}")

    print("\n所有测试样本文件创建完成!")


if __name__ == '__main__':
    main()
